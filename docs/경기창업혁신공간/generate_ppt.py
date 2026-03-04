# -*- coding: utf-8 -*-
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation('붙임2-2. (승급,신규)2차 발표심사  PPT(양식).pptx')
slides = list(prs.slides)

# Colors
BLUE = RGBColor(0x00, 0x70, 0xC0)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
RED_ACCENT = RGBColor(0xE7, 0x4C, 0x3C)
GREEN = RGBColor(0x27, 0xAE, 0x60)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
FONT = '맑은 고딕'
FONT_TITLE = 'HY헤드라인M'


def clear_and_set(text_frame, lines, font_name=FONT, font_size=Pt(12), bold=False, color=DARK):
    while len(text_frame.paragraphs) > 1:
        p = text_frame.paragraphs[-1]._p
        p.getparent().remove(p)
    first_p = text_frame.paragraphs[0]
    for run in list(first_p.runs):
        first_p._p.remove(run._r)

    for idx, line in enumerate(lines):
        if idx == 0:
            para = text_frame.paragraphs[0]
        else:
            para = text_frame.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_before = Pt(1)
        para.space_after = Pt(1)

        if isinstance(line, list):
            for item in line:
                if isinstance(item, tuple):
                    txt, b, sz, c = item
                    run = para.add_run()
                    run.text = txt
                    run.font.name = font_name
                    run.font.size = sz if sz else font_size
                    run.font.bold = b
                    run.font.color.rgb = c if c else color
                else:
                    run = para.add_run()
                    run.text = item
                    run.font.name = font_name
                    run.font.size = font_size
                    run.font.bold = bold
                    run.font.color.rgb = color
        else:
            run = para.add_run()
            run.text = line
            run.font.name = font_name
            run.font.size = font_size
            run.font.bold = bold
            run.font.color.rgb = color


def has_red_text(shape):
    if not shape.has_text_frame:
        return False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                if run.font.color and run.font.color.rgb and str(run.font.color.rgb) == 'FF0000':
                    return True
            except:
                pass
    return False


# ============================================================
# SLIDE 1 - Cover
# ============================================================
slide1 = slides[0]
for shape in slide1.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if '2024' in text:
            clear_and_set(shape.text_frame, ['2026. 03. 26.'],
                          font_name=FONT_TITLE, font_size=Pt(20), color=DARK)
        elif '기업명' in text and '대표' in text:
            clear_and_set(shape.text_frame, [
                [('렌트미 (RentMe)', True, Pt(30), DARK)],
                '',
                [('"세입자의 LinkedIn"', False, Pt(16), BLUE)],
                [('AI 기반 세입자 신뢰 프로필 플랫폼', False, Pt(16), DARK)],
                '',
                [('대표  김다운', True, Pt(18), DARK)],
            ], font_name=FONT_TITLE, font_size=Pt(24), color=DARK)
        elif '기업로고' in text:
            clear_and_set(shape.text_frame, [
                [('RentMe', True, Pt(20), BLUE)],
            ], font_name=FONT_TITLE, font_size=Pt(18), color=BLUE)
        elif has_red_text(shape):
            clear_and_set(shape.text_frame, [
                [('좋은 세입자임을 증명하면 보증금이 달라집니다', False, Pt(11), GRAY)],
            ], font_size=Pt(11))

print('Slide 1 - Cover')

# ============================================================
# SLIDE 3 - 1. 회사소개 → 창업기업우수성 (25점)
# ============================================================
slide3 = slides[2]
for shape in slide3.shapes:
    if has_red_text(shape):
        clear_and_set(shape.text_frame, [
            [('회사 개요', True, Pt(16), BLUE)],
            [('회사명  ', True, Pt(12), None), ('(주)렌트미 (예정) | 예비창업자 | D.N.A - 인공지능', False, Pt(12), None)],
            [('대표자  ', True, Pt(12), None), ('김다운 (1988.11.22) | 풀스택 개발 10년 (2016~현재)', False, Pt(12), None)],
            [('설립예정  ', True, Pt(12), None), ('2026년 6월 (입주 후 3개월 이내 법인 설립)', False, Pt(12), None)],
            [('소재지  ', True, Pt(12), None), ('경기도 하남시 검단산로 239, 창업혁신공간 동부권(하남)', False, Pt(12), None)],
            '',
            [('문제 정의 — 임대차 시장의 "신뢰 인프라" 부재', True, Pt(14), BLUE)],
            [('  1인 가구 804만(36.1%) | 월세 비중 60.2%(역대 최고) | 전세사기 피해 2.5만명', False, Pt(11), None)],
            [('  세입자: ', True, Pt(11), RED_ACCENT), ('"좋은 세입자인데 증명할 방법이 없다"', False, Pt(11), RED_ACCENT)],
            [('  집주인: ', True, Pt(11), RED_ACCENT), ('"월세를 잘 낼지 어떻게 알지?"', False, Pt(11), RED_ACCENT)],
            [('  → 결국 "돈(보증금)"으로만 신뢰를 담보하는 구조', False, Pt(11), None)],
            '',
            [('솔루션 — "세입자의 LinkedIn"', True, Pt(14), BLUE)],
            '  렌트미는 세입자의 신뢰를 정량화하여 "자산"으로 만드는 플랫폼',
            [('  재직/소득/신용 API 자동인증 + 집주인 레퍼런스(한국 최초) + AI 자기소개', False, Pt(11), None)],
            [('  → 신뢰점수(Trust Score) 산출 → 보증금 협상력 확보 + 보증보험 연계', False, Pt(11), None)],
            '',
            [('MVP 99% 완성 — 즉시 서비스 런칭 가능 상태', True, Pt(13), GREEN)],
            [('  (세입자 온보딩, 신뢰점수, 본인인증 기반 자동 스크래핑, AI 자기소개, 집주인 대시보드, 메시징, 매물관리, 보안, 테스트 전 기능 구현)', False, Pt(10), GRAY)],
        ], font_name=FONT, font_size=Pt(12), color=DARK)
        break

print('Slide 3 - Company')

# ============================================================
# SLIDE 4 - 대표자/조직 → 창업기업우수성 (25점)
# ============================================================
slide4 = slides[3]
for shape in slide4.shapes:
    if has_red_text(shape):
        clear_and_set(shape.text_frame, [
            [('대표자 역량 — 렌트미에 최적화된 10년 실무 경험', True, Pt(16), BLUE)],
            '',
            [('김다운  |  ', True, Pt(14), None), ('풀스택 개발 10년 | 방송통신대 컴퓨터과학과 재학', False, Pt(12), None)],
            '',
            [('경력과 렌트미 적용 연결', True, Pt(13), BLUE)],
            [('  헥토데이터(구 쿠콘) 파트장 3년', True, Pt(11), None), (' (2021~2024)', False, Pt(11), GRAY)],
            '    금융 데이터 스크래핑/인증 API 실무 + 팀 관리',
            [('    → 렌트미 재직/소득/신용 자동 스크래핑 인증 시스템 직접 설계/구현', False, Pt(11), GREEN)],
            '',
            [('  IBK 기업은행 파견', True, Pt(11), None), (' (2020~2021)', False, Pt(11), GRAY)],
            '    전자약정 다중인증(MFA) 시스템 구현 (카카오/1원인증)',
            [('    → 렌트미 NICE 본인인증 + SMS OTP 설계', False, Pt(11), GREEN)],
            '',
            [('  에너지마켓플레이스', True, Pt(11), None), (' (2019~2020)', False, Pt(11), GRAY)],
            '    수요-공급 매칭 알고리즘 설계/개발',
            [('    → 렌트미 세입자-집주인 매칭 로직 적용', False, Pt(11), GREEN)],
            '',
            [('  크리에이티브힐', True, Pt(11), None), (' (2018~2019) ', False, Pt(11), GRAY), ('NFT B2C 플랫폼 기획/개발/운영', False, Pt(11), None)],
            [('  지자체 웹서비스', True, Pt(11), None), (' (2016~2018) ', False, Pt(11), GRAY), ('서귀포시/고양시 등 B2G → 공공기관 납품 기반', False, Pt(11), None)],
            '',
            [('핵심: ', True, Pt(12), RED_ACCENT), ('기획→설계→개발→배포→운영 1인 수행 가능 | 퇴근/주말 활용 MVP 99% 단독 개발 완료', False, Pt(12), RED_ACCENT)],
        ], font_name=FONT, font_size=Pt(11), color=DARK)
        break

print('Slide 4 - CEO')

# ============================================================
# SLIDE 5 - 2. 기술성 → 아이템기술성 (30점)
# ============================================================
slide5 = slides[4]
for shape in slide5.shapes:
    if has_red_text(shape):
        clear_and_set(shape.text_frame, [
            [('3대 핵심 기술', True, Pt(15), BLUE)],
            '',
            [('1. 본인인증 기반 자동 스크래핑 인증', True, Pt(12), None), (' — 서류 없이 원클릭, 위변조 불가', False, Pt(11), None)],
            '   사용자 본인인증 후 자체 스크래핑: 국민연금(재직) + 건강보험료(소득) + 신용점수 실시간 조회',
            '   NICE 본인인증(CI/DI) + SMS OTP(6자리) | RSA 암호화 민감정보 보호',
            '',
            [('2. 신뢰점수(Trust Score) 알고리즘', True, Pt(12), None), (' — 0~120점 다차원 정량 평가', False, Pt(11), None)],
            '   프로필(20) + 재직(25) + 소득(25) + 신용(10~20) + 집주인 레퍼런스(+30/-20)',
            '   등급: Excellent(80+) / Good(50~79) / Fair(20~49) / Low(~19)',
            [('   ※ 집주인 레퍼런스(한국 최초): SMS 토큰 기반 설문(월세납부/시설관리/이웃관계/퇴거상태)', False, Pt(10), GRAY)],
            '',
            [('3. AI 자기소개 생성', True, Pt(12), None), (' — OpenAI GPT-4o-mini', False, Pt(11), None)],
            '   프로필 데이터(라이프스타일/거주습관/가족구성) 분석 → 집주인 관점 최적화 소개 자동 생성',
            '',
            [('기술 스택 (전체 직접 구현)', True, Pt(13), BLUE)],
            '  Frontend: Next.js 14, React 18, TypeScript, Tailwind CSS, shadcn/ui',
            '  Backend: Next.js App Router, PostgreSQL, JWT(HS256), bcryptjs',
            '  인프라: Vercel + AWS S3 + Cloudflare R2 + Sentry',
            '  보안: Rate Limiting + CSRF + XSS 방지 + AES-256 + RLS',
            '  테스트: Vitest(단위) + Playwright(E2E)',
            '',
            [('차별점: ', True, Pt(12), RED_ACCENT), ('해외(Naborly/OpenRent)에 유사 서비스 존재하나 한국 시장 특성(전세/월세) 반영 서비스는 전무', False, Pt(11), RED_ACCENT)],
        ], font_name=FONT, font_size=Pt(10), color=DARK)
        break

print('Slide 5 - Tech')

# ============================================================
# SLIDE 6 - 3. 사업성 → 기업구성계획 (20점)
# ============================================================
slide6 = slides[5]
for shape in slide6.shapes:
    if has_red_text(shape):
        clear_and_set(shape.text_frame, [
            [('시장 규모', True, Pt(14), BLUE)],
            '  TAM: 국내 전월세 시장 540조원, 연간 임대차 250만 건',
            '  SAM: 1인 가구 월세 시장 6,000억원 (1인 가구 804만, 월세 120만건/년)',
            '  SOM: 수도권 1인 월세 5% 점유 = 30억원 (3년 목표)',
            [('  글로벌 프롭테크: $182억(2024) → $865억(2034), CAGR 16.9%', False, Pt(10), GRAY)],
            '',
            [('경쟁사 비교', True, Pt(13), BLUE)],
            '  직방/다방: 매물 중심 → 세입자 프로필 없음 | 피터팬: 직거래 → 인증 없음',
            '  안심전세: 건물 위험도만 → 세입자 관점 없음 | 네이버부동산: 시세 정보만',
            [('  → 세입자 "신뢰 프로필" 분야 경쟁자 부재, 선점 기회', True, Pt(11), GREEN)],
            '',
            [('수익 모델', True, Pt(13), BLUE)],
            '  B2C 프리미엄 구독: W9,900/월 (무제한 AI 자기소개, 레퍼런스, 우선 노출)',
            '  집주인 프로필 열람: W3,900/건 | 집주인 구독: W29,900/월 (무제한 열람+매칭)',
            '  중장기: 보증보험 수수료 3~5% (신뢰점수 기반 보증금 50% 절감 연계)',
            '',
            [('Unit Economics', True, Pt(13), BLUE)],
            [('  세입자 CAC W5,000 / LTV W59,400 = LTV:CAC ', False, Pt(11), None), ('12:1', True, Pt(11), GREEN)],
            [('  집주인 CAC W20,000 / LTV W179,400 = LTV:CAC ', False, Pt(11), None), ('9:1', True, Pt(11), GREEN)],
            [('  인프라 BEP: 유료 89명 | 사업 전체 BEP: 유료 약 500명 (2027 하반기)', False, Pt(11), ORANGE)],
            '',
            [('마케팅: ', True, Pt(12), BLUE), ('바이럴 루프', False, Pt(12), None)],
            '  세입자 가입→프로필 완성→링크 공유→집주인 유입→세입자 확대 선순환',
            '  온라인: 부동산 커뮤니티(40%) + 블로그SEO/유튜브(30%) + SNS(20%)',
            '  B2B: 셰어하우스(우주/맹그로브) B2B 제휴 → 대학가 → 동부권 중개업소 확대',
        ], font_name=FONT, font_size=Pt(10), color=DARK)
        break

print('Slide 6 - Business')

# ============================================================
# SLIDE 7 - 4. 성장성 → 투자유치계획 (25점)
# ============================================================
slide7 = slides[6]
for shape in slide7.shapes:
    if has_red_text(shape):
        clear_and_set(shape.text_frame, [
            [('단계별 성장 로드맵', True, Pt(15), BLUE)],
            '',
            [('Phase 1 | 2026 하반기 — MVP 출시 + PMF 검증', True, Pt(12), None)],
            '  정식 런칭(7월) | 셰어하우스 B2B 제휴 | MAU 3,000명 | 유료전환 3%+',
            '  B2C 프리미엄 구독 출시 | 매출 500만원 | 고용 1명(대표)',
            '',
            [('Phase 2 | 2027년 — B2B 확장 + 투자유치', True, Pt(12), None)],
            '  B2B SaaS 출시(중개업소 100곳 제휴) | MAU 20,000명 | 유료전환 8%+',
            [('  Seed 투자 2~3억 (VC/AC) | ', False, Pt(11), None), ('매출 1.5억원 | 고용 3명', True, Pt(11), None)],
            '  팀 확장: 프론트엔드 개발자 1명 + 마케터 1명 채용',
            '',
            [('Phase 3 | 2028년~ — 플랫폼 확장 + BEP', True, Pt(12), None)],
            '  공공기관 연계(LH/SH 세입자 스크리닝) | 모바일 앱(React Native) 출시',
            '  양방향 신뢰 플랫폼(집주인 인증 도입) | 보증보험사 제휴 본격화',
            [('  Series A 10억 | ', False, Pt(11), None), ('매출 8억원 (BEP 달성) | 고용 6명', True, Pt(11), None)],
            '',
            [('자금조달 계획', True, Pt(13), BLUE)],
            '  예비창업패키지(최대 1억원) 신청 중 | G스타 오디션 경기창업공모 신청 중',
            '  1인 개발/운영 + 무상 입주 → 월 고정비 약 30만원(서버/API)',
            '  정부지원금 → 초기 사업화 자금(마케팅/API/디자인) | 매출 발생 후 투자유치',
            '',
            [('KPI 목표', True, Pt(13), BLUE)],
            '  M3: 가입 1,000 / M6: MAU 3,000 + 유료 90명 / M12: MAU 10,000 + B2B 50곳',
            '  M24: MAU 30,000 + 유료전환 8% + B2B 200곳',
        ], font_name=FONT, font_size=Pt(10), color=DARK)
        break

print('Slide 7 - Growth')

# ============================================================
# SLIDE 8 - 5. 지식재산권, 인증, 수상
# ============================================================
slide8 = slides[7]
for shape in slide8.shapes:
    if has_red_text(shape):
        clear_and_set(shape.text_frame, [
            [('지식재산권 계획', True, Pt(14), BLUE)],
            '  "렌트미(RentMe)" 서비스 상표 출원 예정',
            '  다층 신뢰 스코어링 알고리즘 특허 출원 예정',
            '  AI 기반 세입자 자기소개 생성 방법 특허 출원 예정',
            '',
            [('기술 인증 현황 (MVP 구현 완료)', True, Pt(14), BLUE)],
            '  자체 스크래핑 인증 구현 완료 — 본인인증 후 재직/소득/신용 자동 수집',
            '  NICE API 연동 완료 — 본인인증 CI/DI 발급',
            '  OpenAI GPT-4o-mini 연동 완료 — AI 자기소개 생성',
            '  소셜 로그인 3종 연동 완료 — 카카오/네이버/구글 OAuth',
            '  보안 시스템 구현 완료 — Rate Limiting, CSRF, XSS, AES-256, RLS',
            '  테스트 환경 구축 완료 — Vitest(단위) + Playwright(E2E)',
            '',
            [('창업지원사업 현황', True, Pt(14), BLUE)],
            '  2026 예비창업패키지 (중소벤처기업부, 최대 1억원) — 신청 중',
            '  2026 경기창업공모 G스타 오디션 (예비/초기리그) — 신청 중',
            '  2026 경기창업혁신공간 동부권(하남) 틔움 입주 — 신청 중',
            '',
            [('가점 해당', True, Pt(14), BLUE)],
            [('  신산업 창업분야 (D.N.A - 인공지능) — ', False, Pt(12), None), ('가점 1점', True, Pt(12), GREEN)],
        ], font_name=FONT, font_size=Pt(11), color=DARK)
        break

print('Slide 8 - IP/Cert')

# ============================================================
# SLIDE 9 - 6. 기타사항
# ============================================================
slide9 = slides[8]
for shape in slide9.shapes:
    if has_red_text(shape):
        clear_and_set(shape.text_frame, [
            [('입주 공간 필요성', True, Pt(14), BLUE)],
            '  예비창업자 → 사업자등록을 위한 사업장 주소 필요',
            '  입주 확정 시 퇴사 → 풀타임 전환 → 즉시 법인설립 + 서비스 런칭',
            '  동부권(하남/광주/성남) 부동산 중개업소 B2B 영업 지리적 접근성 우수',
            '',
            [('공간 활용 계획', True, Pt(13), BLUE)],
            '  제품 개발/운영: 서비스 고도화, 모바일 앱 개발, 상시 운영',
            '  사업장 등록: 입주 후 3개월 이내 법인 설립 + 사업자등록(본사 주소)',
            '  B2B 미팅: 동부권 부동산 중개업소 영업 거점',
            '  투자 IR: 센터 내 회의실 활용, VC/AC 투자 미팅 진행',
            '',
            [('입주 후 추진 일정', True, Pt(13), BLUE)],
            '  2026.04~06  퇴사 및 법인 설립, 풀타임 전환',
            '  2026.06~07  서비스 정식 런칭 (웹)',
            '  2026.07~09  셰어하우스 B2B 제휴, 초기 사용자 1,000명 확보',
            '  2026.09~12  B2C 프리미엄 출시, MAU 3,000명 달성',
            '  2027.01~     모바일 앱 출시, B2B SaaS, Seed 투자 유치',
            '',
            [('기대효과', True, Pt(13), BLUE)],
            '  무상 입주 → 초기 고정비 절감(임대료 0원), 사업화 자금에 집중',
            '  입주기업 네트워킹/멘토링/오픈이노베이션 → BM 고도화 + 투자 준비',
            '  경기도 창업지원사업/전세피해지원센터 등 지역 자원 연계 활용',
        ], font_name=FONT, font_size=Pt(11), color=DARK)
        break

print('Slide 9 - Others')

# Save
output_path = '붙임2-2_발표심사PPT_렌트미_김다운.pptx'
prs.save(output_path)
print(f'\nSaved: {output_path}')
print('All slides completed!')
