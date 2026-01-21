from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 프레젠테이션 생성
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# 색상 정의
PRIMARY = RGBColor(41, 98, 255)
PRIMARY_DARK = RGBColor(25, 60, 180)
ACCENT = RGBColor(255, 87, 34)
ACCENT2 = RGBColor(76, 175, 80)
DARK = RGBColor(33, 33, 33)
GRAY = RGBColor(117, 117, 117)
LIGHT_GRAY = RGBColor(245, 245, 245)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(244, 67, 54)
YELLOW = RGBColor(255, 193, 7)


# ========== 헬퍼 함수 ==========

def add_background(slide, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_header_bar(slide, height=1.0):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    return bar


def add_text_box(slide, left, top, width, height, text, font_size=24, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_bullet_text(slide, left, top, width, height, items, font_size=22, color=DARK, bullet="•"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet} {item}" if bullet else item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(12)
    return box


def add_icon_card(slide, left, top, width, height, icon, title, content):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(230, 230, 230)
    add_text_box(slide, left + 0.3, top + 0.3, width - 0.6, 0.6, icon, font_size=36, align=PP_ALIGN.CENTER)
    add_text_box(slide, left + 0.3, top + 1.0, width - 0.6, 0.5, title, font_size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide, left + 0.3, top + 1.5, width - 0.6, height - 1.8, content, font_size=14, color=GRAY, align=PP_ALIGN.CENTER)


def add_stat_box(slide, left, top, number, label, color=PRIMARY):
    add_text_box(slide, left, top, 3, 0.8, number, font_size=48, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + 0.9, 3, 0.5, label, font_size=16, color=GRAY, align=PP_ALIGN.CENTER)


# ========== 슬라이드 생성 함수 ==========

def slide_01_cover(prs):
    """1. 표지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, PRIMARY)
    add_text_box(slide, 1, 2.5, 14, 1.2, "Fryndo AR Companion",
                 font_size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 4.0, 14, 0.8, "언어가 달라도 함께 여행하는 AR 동행 서비스",
                 font_size=32, color=WHITE, align=PP_ALIGN.CENTER)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(5.2), Inches(4), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = WHITE
    line.line.fill.background()
    add_text_box(slide, 1, 5.8, 14, 0.5, "2026 예비창업패키지 IR",
                 font_size=20, color=WHITE, align=PP_ALIGN.CENTER)
    # [수정 필요] 발표자 이름
    add_text_box(slide, 1, 6.4, 14, 0.5, "발표자: [이름]",
                 font_size=18, color=WHITE, align=PP_ALIGN.CENTER)


def slide_02_toc(prs):
    """2. 목차"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, 1.2)
    add_text_box(slide, 0.8, 0.35, 14, 0.7, "목차", font_size=40, bold=True, color=WHITE)

    items_left = [
        "1. 문제 정의",
        "2. 솔루션 & 핵심 기능",
        "3. 사용자 시나리오",
        "4. 데모 시연",
        "5. 기존 MVP 현황",
    ]
    items_right = [
        "6. 시장 분석",
        "7. 경쟁력 & 비즈니스 모델",
        "8. 팀 & 로드맵",
        "9. 자금 계획",
        "10. 리스크 대응",
    ]

    for i, item in enumerate(items_left):
        add_text_box(slide, 2, 2.0 + i * 1.1, 5, 0.8, item, font_size=28, color=DARK)
    for i, item in enumerate(items_right):
        add_text_box(slide, 9, 2.0 + i * 1.1, 5, 0.8, item, font_size=28, color=DARK)


def slide_03_problem1(prs):
    """3. 문제 정의 - 언어장벽"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, 1.2)
    add_text_box(slide, 0.8, 0.35, 14, 0.7, "문제 정의 ①  언어 장벽", font_size=40, bold=True, color=WHITE)
    add_text_box(slide, 0.8, 1.6, 14, 0.6, "해외여행 중 외국인과 함께하게 되면...",
                 font_size=24, color=GRAY)

    add_icon_card(slide, 0.8, 2.4, 4.5, 2.8, "🗣️", "대화 단절",
                  "말이 안 통해서\n어색한 침묵이 계속됨")
    add_icon_card(slide, 5.75, 2.4, 4.5, 2.8, "📱", "번역 앱 한계",
                  "1:1만 지원\n3명 이상 그룹은 불가능")
    add_icon_card(slide, 10.7, 2.4, 4.5, 2.8, "⏱️", "시간 낭비",
                  "일일이 번역하느라\n여행을 제대로 못 즐김")

    # 인용 (78% 통계 삭제, 실제 경험담으로 대체)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.8), Inches(12), Inches(1.8))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.fill.background()
    add_text_box(slide, 2.5, 6.0, 11, 0.5, "💬 실제 여행자 후기", font_size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide, 2.5, 6.5, 11, 1,
                 '"일본인, 미국인이랑 같이 여행했는데 번역 앱으로 한 명씩 말해야 해서 너무 불편했어요.\n그룹 대화가 되는 앱이 있으면 좋겠다고 생각했어요."',
                 font_size=16, color=GRAY, align=PP_ALIGN.CENTER)


def slide_04_problem2(prs):
    """4. 문제 정의 - 만남/길찾기"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, 1.2)
    add_text_box(slide, 0.8, 0.35, 14, 0.7, "문제 정의 ②  만남 & 길찾기", font_size=40, bold=True, color=WHITE)

    add_text_box(slide, 0.8, 1.8, 6.5, 0.5, "🤝 동행자 만남", font_size=28, bold=True, color=DARK)
    problems_left = [
        "약속 장소에 도착했는데 서로 못 찾음",
        "사람 많은 곳에서 전화해도 설명 어려움",
        "외모 특징 설명하기 난감 (외국인끼리)",
        "30분 이상 헤매다 지침"
    ]
    add_bullet_text(slide, 0.8, 2.5, 6.5, 3, problems_left, font_size=20)

    add_text_box(slide, 8.5, 1.8, 6.5, 0.5, "🗺️ 낯선 도시 길찾기", font_size=28, bold=True, color=DARK)
    problems_right = [
        "지도 앱 보느라 고개 숙이고 걸음",
        "주변 풍경을 제대로 못 봄",
        "동행자들이 뿔뿔이 흩어짐",
        "여행의 즐거움 반감"
    ]
    add_bullet_text(slide, 8.5, 2.5, 6.5, 3, problems_right, font_size=20)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.8), Inches(14), Inches(2))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 243, 224)
    box.line.fill.background()
    add_text_box(slide, 1.5, 6.0, 13, 0.6, "⚠️ 기존 솔루션의 한계", font_size=24, bold=True, color=ACCENT)
    add_text_box(slide, 1.5, 6.6, 13, 1,
                 "번역 앱: 1:1만 지원  |  지도 앱: 화면 계속 봐야 함  |  여행 앱: 기능이 분리됨",
                 font_size=20, color=DARK, align=PP_ALIGN.CENTER)


def slide_05_solution(prs):
    """5. 솔루션 소개"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, PRIMARY)
    add_text_box(slide, 1, 1.5, 14, 0.8, "Solution", font_size=24, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 2.2, 14, 1.2, "Fryndo AR Companion",
                 font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 3.8, 14, 0.8, "여행 동행의 모든 문제를 AR로 해결합니다",
                 font_size=28, color=WHITE, align=PP_ALIGN.CENTER)

    features = [
        ("🗣️", "그룹 번역", "3명 이상\n실시간 번역"),
        ("📍", "AR 찾기", "동행자 위치\nAR 표시"),
        ("🧭", "AR 길안내", "화살표로\n직관적 안내"),
    ]

    for i, (icon, title, desc) in enumerate(features):
        x = 2 + i * 4.5
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5), Inches(3.8), Inches(3))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.fill.background()
        add_text_box(slide, x, 5.2, 3.8, 0.6, icon, font_size=40, align=PP_ALIGN.CENTER)
        add_text_box(slide, x, 5.9, 3.8, 0.5, title, font_size=22, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
        add_text_box(slide, x, 6.5, 3.8, 1.2, desc, font_size=16, color=GRAY, align=PP_ALIGN.CENTER)


def slide_06_feature1(prs):
    """6. 핵심 기능 1 - 그룹 번역"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, 1.2)
    add_text_box(slide, 0.8, 0.35, 14, 0.7, "핵심 기능 ①  실시간 그룹 번역", font_size=40, bold=True, color=WHITE)

    add_text_box(slide, 0.8, 1.8, 7, 0.6, "3명 이상 다국어 동시 번역", font_size=28, bold=True, color=PRIMARY)
    features = [
        "음성 인식 → 텍스트 → 번역 → 자막 표시",
        "한국어, 영어, 일본어, 중국어 등 20개 언어",
        "그룹 채팅방에서 각자 모국어로 표시",
        "음성 + 텍스트 동시 지원",
        "오프라인 모드 지원 (향후)"
    ]
    add_bullet_text(slide, 0.8, 2.6, 7, 3.5, features, font_size=20)

    add_text_box(slide, 0.8, 6.0, 7, 0.5, "사용 기술", font_size=18, bold=True, color=GRAY)
    add_text_box(slide, 0.8, 6.5, 7, 0.8, "음성인식: Whisper API  |  번역: DeepL API", font_size=16, color=GRAY)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.8), Inches(6.5), Inches(6))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.color.rgb = RGBColor(200, 200, 200)
    add_text_box(slide, 8.8, 2.0, 6, 0.5, "💬 그룹 채팅 시뮬레이션", font_size=18, bold=True, color=DARK)

    messages = [
        ("🇺🇸 Mike", "Let's grab some ramen!", "라멘 먹으러 가자!", 2.7),
        ("🇯🇵 Yuki", "いいね！この近くにおいしい店がある", "좋아! 이 근처에 맛있는 집 있어", 4.0),
        ("🇰🇷 나", "좋아요! 어디로 가면 돼요?", "Where should we go?", 5.3),
    ]
    for name, original, translated, y in messages:
        add_text_box(slide, 8.8, y, 6, 0.3, name, font_size=14, bold=True, color=DARK)
        add_text_box(slide, 8.8, y + 0.35, 6, 0.3, original, font_size=16, color=DARK)
        add_text_box(slide, 8.8, y + 0.7, 6, 0.3, f"→ {translated}", font_size=14, color=PRIMARY)


def slide_07_feature2(prs):
    """7. 핵심 기능 2 - AR 동행자 찾기"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, 1.2)
    add_text_box(slide, 0.8, 0.35, 14, 0.7, "핵심 기능 ②  AR 동행자 찾기", font_size=40, bold=True, color=WHITE)

    add_text_box(slide, 0.8, 1.8, 7, 0.6, "카메라로 주변 스캔, 동행자 위치 AR 표시", font_size=24, bold=True, color=PRIMARY)
    features = [
        "스마트폰 카메라로 주변 스캔",
        "동행자 프로필 + 거리 AR 오버레이",
        '"3시 방향 15m" 직관적 표시',
        "GPS + 블루투스 기반 정확한 위치",
        "프로필 사진으로 쉽게 구분"
    ]
    add_bullet_text(slide, 0.8, 2.6, 7, 3.5, features, font_size=20)

    add_text_box(slide, 0.8, 6.0, 7, 0.5, "사용 흐름", font_size=18, bold=True, color=GRAY)
    add_text_box(slide, 0.8, 6.5, 7, 0.8, "약속 장소 도착 → 앱 실행 → 카메라 스캔 → 동행자 발견!", font_size=16, color=GRAY)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.8), Inches(6.5), Inches(6))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(50, 50, 50)
    box.line.color.rgb = RGBColor(100, 100, 100)
    add_text_box(slide, 8.8, 2.0, 6, 0.4, "📷 AR 카메라 뷰", font_size=16, color=WHITE)

    marker = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(3.5), Inches(2.5), Inches(1.5))
    marker.fill.solid()
    marker.fill.fore_color.rgb = PRIMARY
    marker.line.fill.background()
    add_text_box(slide, 10.5, 3.6, 2.5, 0.4, "🇺🇸 Mike", font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 10.5, 4.1, 2.5, 0.4, "15m →", font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 10.5, 4.6, 2.5, 0.3, "3시 방향", font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, 8.8, 6.8, 6, 0.5, "※ 실제 카메라 영상 위에 AR 정보가 오버레이됩니다",
                 font_size=12, color=GRAY, align=PP_ALIGN.CENTER)


def slide_08_feature3(prs):
    """8. 핵심 기능 3 - AR 길안내"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, 1.2)
    add_text_box(slide, 0.8, 0.35, 14, 0.7, "핵심 기능 ③  AR 길안내", font_size=40, bold=True, color=WHITE)

    add_text_box(slide, 0.8, 1.8, 7, 0.6, "화면에 화살표, 지도 안 봐도 OK", font_size=24, bold=True, color=PRIMARY)
    features = [
        "카메라 화면에 방향 화살표 오버레이",
        "목적지까지 남은 거리 실시간 표시",
        "동행자 위치도 함께 표시",
        "그룹 전체가 같은 목적지 공유",
        "음성 안내 병행 가능"
    ]
    add_bullet_text(slide, 0.8, 2.6, 7, 3.5, features, font_size=20)

    add_text_box(slide, 0.8, 6.0, 7, 0.5, "기존 지도 앱 대비 장점", font_size=18, bold=True, color=GRAY)
    add_text_box(slide, 0.8, 6.5, 7, 0.8, "고개 안 숙여도 됨 → 주변 풍경 즐기면서 이동 가능", font_size=16, color=ACCENT2)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.8), Inches(6.5), Inches(6))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(60, 60, 60)
    box.line.color.rgb = RGBColor(100, 100, 100)
    add_text_box(slide, 8.8, 2.0, 6, 0.4, "🧭 AR 내비게이션", font_size=16, color=WHITE)
    add_text_box(slide, 10, 3.5, 3, 1, "⬆️", font_size=80, color=ACCENT2, align=PP_ALIGN.CENTER)

    dest_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), Inches(5.5), Inches(4), Inches(1.2))
    dest_box.fill.solid()
    dest_box.fill.fore_color.rgb = WHITE
    dest_box.line.fill.background()
    add_text_box(slide, 9.5, 5.6, 4, 0.4, "🍜 이치란 라멘", font_size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide, 9.5, 6.1, 4, 0.4, "직진 200m", font_size=24, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)


def slide_09_scenario(prs):
    """9. 사용자 시나리오"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, 1.2)
    add_text_box(slide, 0.8, 0.35, 14, 0.7, "사용자 시나리오: 도쿄 여행", font_size=40, bold=True, color=WHITE)

    steps = [
        ("1", "매칭", "한국인 A, 미국인 B,\n일본인 C가 Fryndo로\n동행 매칭", 1.0),
        ("2", "만남", "시부야역 하치코 앞\n약속, AR로 서로\n쉽게 찾음", 4.0),
        ("3", "대화", "3개 국어 실시간\n번역으로 자연스럽게\n대화", 7.0),
        ("4", "이동", "AR 길안내로\n라멘집까지\n함께 이동", 10.0),
        ("5", "여행", "언어 장벽 없이\n즐거운 여행\n완성!", 13.0),
    ]

    for num, title, desc, x in steps:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.2), Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY
        circle.line.fill.background()
        add_text_box(slide, x, 2.3, 0.7, 0.5, num, font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, x - 0.3, 3.1, 1.5, 0.5, title, font_size=20, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text_box(slide, x - 0.5, 3.7, 2, 1.5, desc, font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

    for i in range(4):
        add_text_box(slide, 2.0 + i * 3, 2.35, 1.5, 0.5, "→", font_size=24, color=GRAY, align=PP_ALIGN.CENTER)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.8), Inches(14), Inches(2))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(232, 245, 253)
    box.line.fill.background()
    add_text_box(slide, 1.5, 6.0, 13, 0.5, "💡 핵심 가치", font_size=22, bold=True, color=PRIMARY)
    add_text_box(slide, 1.5, 6.6, 13, 1,
                 "처음 만난 외국인과도 오랜 친구처럼 여행할 수 있는 경험",
                 font_size=24, color=DARK, align=PP_ALIGN.CENTER)


def slide_10_demo(prs):
    """10. 데모 시연"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK)
    add_text_box(slide, 1, 2, 14, 0.6, "DEMO", font_size=32, color=GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 2.8, 14, 1.2, "🎬 데모 영상", font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(4.5), Inches(8), Inches(2))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(60, 60, 60)
    box.line.color.rgb = GRAY
    box.line.width = Pt(2)
    add_text_box(slide, 4, 5.0, 8, 0.5, "▶️  45초 데모 영상 삽입", font_size=24, color=GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, 4, 5.7, 8, 0.5, "그룹 번역 → AR 동행자 찾기 → AR 길안내", font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

    add_text_box(slide, 1, 7.5, 14, 0.5, "※ 현재 개발 중이며, 목업 기반 데모입니다",
                 font_size=14, color=GRAY, align=PP_ALIGN.CENTER)


def slide_11_mvp(prs):
    """11. 기존 Fryndo MVP 현황 (NEW)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, 1.2)
    add_text_box(slide, 0.8, 0.35, 14, 0.7, "기존 Fryndo MVP 현황", font_size=40, bold=True, color=WHITE)

    # 왼쪽: MVP 설명
    add_text_box(slide, 0.8, 1.8, 7, 0.6, "이미 검증된 동행 매칭 MVP 보유", font_size=24, bold=True, color=PRIMARY)

    mvp_features = [
        "여행 동행자 매칭 기능 구현 완료",
        "실시간 채팅 기능",
        "여행 일정 공유 기능",
        "React Native 크로스플랫폼 앱",
        "Node.js 백엔드 (Koyeb 배포)"
    ]
    add_bullet_text(slide, 0.8, 2.6, 7, 3, mvp_features, font_size=20)

    # MVP 현황 박스
    status_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.5), Inches(7), Inches(2.2))
    status_box.fill.solid()
    status_box.fill.fore_color.rgb = RGBColor(232, 245, 233)
    status_box.line.fill.background()

    add_text_box(slide, 1.2, 5.7, 6, 0.5, "✅ MVP 개발 현황", font_size=20, bold=True, color=ACCENT2)
    add_text_box(slide, 1.2, 6.3, 6, 1.2,
                 "• 핵심 기능 개발 완료\n• 내부 테스트 진행 중\n• AR 기능 확장 준비 단계",
                 font_size=16, color=DARK)

    # 오른쪽: 앱 화면 플레이스홀더
    add_text_box(slide, 8.5, 1.8, 6.5, 0.5, "📱 Fryndo MVP 화면", font_size=20, bold=True, color=DARK)

    # 앱 화면 목업 3개
    screens = ["동행 매칭", "채팅", "일정 공유"]
    for i, screen in enumerate(screens):
        x = 8.5 + i * 2.2
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.4), Inches(2), Inches(4))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = RGBColor(200, 200, 200)
        add_text_box(slide, x, 4.0, 2, 0.5, screen, font_size=14, color=GRAY, align=PP_ALIGN.CENTER)
        add_text_box(slide, x, 4.5, 2, 0.5, "[스크린샷]", font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

    add_text_box(slide, 8.5, 6.8, 6.5, 0.5, "※ 실제 앱 스크린샷으로 교체 필요",
                 font_size=12, color=GRAY, align=PP_ALIGN.CENTER)


def slide_12_tech(prs):
    """12. 기술 구조"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, 1.2)
    add_text_box(slide, 0.8, 0.35, 14, 0.7, "기술 구조", font_size=40, bold=True, color=WHITE)

    # 모바일 앱
    app_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(1.8), Inches(4), Inches(1.2))
    app_box.fill.solid()
    app_box.fill.fore_color.rgb = PRIMARY
    app_box.line.fill.background()
    add_text_box(slide, 6, 2.0, 4, 0.4, "📱 Fryndo AR Companion", font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 6, 2.5, 4, 0.4, "React Native + AR", font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, 7.5, 3.1, 1, 0.5, "↓", font_size=24, color=GRAY, align=PP_ALIGN.CENTER)

    # 서버
    server_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(3.7), Inches(4), Inches(1.2))
    server_box.fill.solid()
    server_box.fill.fore_color.rgb = ACCENT2
    server_box.line.fill.background()
    add_text_box(slide, 6, 3.9, 4, 0.4, "🖥️ Fryndo Server", font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 6, 4.4, 4, 0.4, "Node.js (Koyeb)", font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, 4.5, 4.1, 1, 0.5, "←", font_size=24, color=GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, 10.5, 4.1, 1, 0.5, "→", font_size=24, color=GRAY, align=PP_ALIGN.CENTER)

    # 외부 API들
    apis = [
        ("🗣️ 번역 API", "DeepL", 1.5, 3.7),
        ("🎤 음성 API", "Whisper", 1.5, 5.3),
        ("🗺️ 지도 API", "Google Maps", 11.5, 3.7),
        ("📍 위치 API", "GPS/BLE", 11.5, 5.3),
    ]
    for title, tech, x, y in apis:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.8), Inches(1))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = RGBColor(200, 200, 200)
        add_text_box(slide, x, y + 0.1, 2.8, 0.4, title, font_size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text_box(slide, x, y + 0.5, 2.8, 0.4, tech, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

    add_text_box(slide, 0.8, 7.0, 14, 0.5,
                 "Frontend: React Native  |  AR: ARCore/ARKit  |  Backend: Node.js  |  Infra: Koyeb",
                 font_size=16, color=GRAY, align=PP_ALIGN.CENTER)


def slide_13_market(prs):
    """13. 시장 분석 (출처 추가)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, 1.2)
    add_text_box(slide, 0.8, 0.35, 14, 0.7, "시장 분석", font_size=40, bold=True, color=WHITE)

    # TAM/SAM/SOM
    add_text_box(slide, 0.8, 1.8, 4.5, 0.5, "🌍 TAM (전체 시장)", font_size=22, bold=True, color=DARK)
    add_text_box(slide, 0.8, 2.4, 4.5, 0.6, "글로벌 해외여행자", font_size=18, color=GRAY)
    add_stat_box(slide, 0.8, 2.9, "15억 명", "연간", PRIMARY)
    add_text_box(slide, 0.8, 4.3, 4.5, 0.4, "출처: UNWTO 2024", font_size=12, color=GRAY)

    add_text_box(slide, 5.75, 1.8, 4.5, 0.5, "🇰🇷 SAM (유효 시장)", font_size=22, bold=True, color=DARK)
    add_text_box(slide, 5.75, 2.4, 4.5, 0.6, "한국 출국 + 방한 외국인", font_size=18, color=GRAY)
    add_stat_box(slide, 5.75, 2.9, "4,500만", "연간", PRIMARY)
    add_text_box(slide, 5.75, 4.3, 4.5, 0.4, "출처: 한국관광공사 2024", font_size=12, color=GRAY)

    add_text_box(slide, 10.7, 1.8, 4.5, 0.5, "🎯 SOM (목표 시장)", font_size=22, bold=True, color=DARK)
    add_text_box(slide, 10.7, 2.4, 4.5, 0.6, "다국적 그룹 여행자", font_size=18, color=GRAY)
    add_stat_box(slide, 10.7, 2.9, "100만 명", "연간 추정", ACCENT)
    add_text_box(slide, 10.7, 4.3, 4.5, 0.4, "자체 추정치", font_size=12, color=GRAY)

    # 시장 트렌드
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.0), Inches(14.4), Inches(2.8))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.fill.background()

    add_text_box(slide, 1.3, 5.2, 13, 0.5, "📈 시장 트렌드", font_size=20, bold=True, color=DARK)

    trends = [
        ("해외여행 폭발적 회복", "2024년 한국인 출국 2,800만 명\n(한국관광공사)"),
        ("다국적 여행 증가", "SNS 기반 글로벌 친구 만들기\n트렌드 확산"),
        ("AR 기술 대중화", "iOS/Android 기기 90%+\nAR 지원 (Apple/Google)"),
    ]

    for i, (title, desc) in enumerate(trends):
        x = 1.3 + i * 4.8
        add_text_box(slide, x, 5.8, 4.5, 0.4, f"✓ {title}", font_size=16, bold=True, color=PRIMARY)
        add_text_box(slide, x, 6.3, 4.5, 1, desc, font_size=13, color=GRAY)


def slide_14_competitor(prs):
    """14. 경쟁력 분석 (확대)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, 1.2)
    add_text_box(slide, 0.8, 0.35, 14, 0.7, "경쟁력 분석", font_size=40, bold=True, color=WHITE)

    # 비교 테이블
    headers = ["기능", "Google\n번역", "Papago", "트리플", "Fryndo\nAR"]
    x_positions = [0.8, 3.5, 5.8, 8.1, 10.8]

    for x, header in zip(x_positions, headers):
        color = PRIMARY if "Fryndo" in header else DARK
        bold = "Fryndo" in header
        add_text_box(slide, x, 1.7, 2.5, 0.7, header, font_size=16, bold=bold, color=color, align=PP_ALIGN.CENTER)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.5), Inches(14.4), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = GRAY
    line.line.fill.background()

    comparisons = [
        ("그룹 번역 (3명+)", "❌", "❌", "❌", "✅"),
        ("실시간 음성", "✅", "✅", "❌", "✅"),
        ("AR 기능", "❌", "❌", "❌", "✅"),
        ("동행 매칭", "❌", "❌", "❌", "✅"),
        ("여행 일정 공유", "❌", "❌", "✅", "✅"),
        ("위치 공유", "❌", "❌", "△", "✅"),
    ]

    for i, row in enumerate(comparisons):
        y = 2.7 + i * 0.6
        for j, val in enumerate(row):
            x = x_positions[j]
            if j == 0:
                add_text_box(slide, x, y, 2.5, 0.5, val, font_size=14, color=DARK, align=PP_ALIGN.CENTER)
            else:
                color = ACCENT2 if val == "✅" else (YELLOW if val == "△" else DARK)
                add_text_box(slide, x, y, 2.5, 0.5, val, font_size=16, color=color, align=PP_ALIGN.CENTER)

    # 핵심 차별점
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.5), Inches(14.4), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(232, 245, 233)
    box.line.fill.background()

    add_text_box(slide, 1.3, 6.7, 13.5, 0.9,
                 "🏆 핵심 차별점: 번역 앱도, 여행 앱도 아닌 '다국적 동행 특화 플랫폼' + AR 기술의 유일한 결합",
                 font_size=18, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)


def slide_15_business(prs):
    """15. 비즈니스 모델"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, 1.2)
    add_text_box(slide, 0.8, 0.35, 14, 0.7, "비즈니스 모델", font_size=40, bold=True, color=WHITE)

    # 무료 티어
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.8), Inches(4.5), Inches(3.2))
    box1.fill.solid()
    box1.fill.fore_color.rgb = LIGHT_GRAY
    box1.line.color.rgb = RGBColor(200, 200, 200)
    add_text_box(slide, 1, 1.9, 4.5, 0.5, "Free", font_size=24, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 2.4, 4.5, 0.6, "₩0", font_size=36, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_bullet_text(slide, 1.3, 3.1, 4, 1.8, ["월 30분 번역", "기본 AR 기능", "동행 매칭 3회/월"], font_size=15, color=GRAY, bullet="•")

    # Plus 티어
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(1.8), Inches(4.5), Inches(3.2))
    box2.fill.solid()
    box2.fill.fore_color.rgb = PRIMARY
    box2.line.fill.background()
    add_text_box(slide, 6, 1.9, 4.5, 0.5, "Plus", font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 6, 2.4, 4.5, 0.6, "₩9,900/월", font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_bullet_text(slide, 6.3, 3.1, 4, 1.8, ["무제한 번역", "전체 AR 기능", "무제한 매칭"], font_size=15, color=WHITE, bullet="✓")

    # B2B
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11), Inches(1.8), Inches(4), Inches(3.2))
    box3.fill.solid()
    box3.fill.fore_color.rgb = ACCENT
    box3.line.fill.background()
    add_text_box(slide, 11, 1.9, 4, 0.5, "B2B (향후)", font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 11, 2.4, 4, 0.6, "별도 협의", font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_bullet_text(slide, 11.3, 3.1, 3.5, 1.8, ["여행사 제휴", "관광청 협업"], font_size=15, color=WHITE, bullet="•")

    # 1년차 목표
    add_text_box(slide, 0.8, 5.5, 14, 0.5, "📊 1년차 목표", font_size=22, bold=True, color=DARK)

    goals = [
        ("앱 다운로드", "10,000회"),
        ("활성 사용자", "1,000명"),
        ("유료 전환", "100명 (10%)"),
        ("월 매출", "₩990,000"),
    ]
    for i, (label, value) in enumerate(goals):
        x = 1 + i * 3.7
        add_text_box(slide, x, 6.1, 3.5, 0.4, label, font_size=16, color=GRAY, align=PP_ALIGN.CENTER)
        add_text_box(slide, x, 6.6, 3.5, 0.5, value, font_size=24, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)


def slide_16_team(prs):
    """16. 팀 소개 (경력 수정)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, 1.2)
    add_text_box(slide, 0.8, 0.35, 14, 0.7, "팀 소개", font_size=40, bold=True, color=WHITE)

    # 대표 프로필
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.8), Inches(7), Inches(4))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.fill.background()

    add_text_box(slide, 1.5, 2.0, 6, 0.5, "👤 대표 / 개발", font_size=24, bold=True, color=DARK)

    # [수정] 5~10년차로 변경
    skills = [
        "React Native / Node.js 풀스택 개발",
        "Fryndo MVP 단독 개발 완료",
        "5~10년차 소프트웨어 엔지니어",
        "스타트업 개발 경험"
    ]
    add_bullet_text(slide, 1.5, 2.7, 6, 3, skills, font_size=18, color=DARK)

    # 협업 파트너
    add_text_box(slide, 8.5, 1.8, 6.5, 0.5, "🤝 협업 파트너", font_size=22, bold=True, color=DARK)

    partners = [
        ("AR 개발", "전문 외주 협업 예정\n(ARCore/ARKit 경험자)"),
        ("UI/UX", "크몽/숨고 활용\n(여행 앱 포트폴리오)"),
        ("QA", "베타 테스트 시\n프리랜서 활용"),
    ]
    for i, (role, desc) in enumerate(partners):
        y = 2.5 + i * 1.5
        add_text_box(slide, 8.5, y, 2.5, 0.4, role, font_size=16, bold=True, color=PRIMARY)
        add_text_box(slide, 11, y, 4, 0.9, desc, font_size=14, color=GRAY)

    # 강점
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(6.2), Inches(14), Inches(1.5))
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(232, 245, 253)
    box2.line.fill.background()
    add_text_box(slide, 1.5, 6.4, 13, 0.5, "💪 팀 강점", font_size=20, bold=True, color=PRIMARY)
    add_text_box(slide, 1.5, 6.9, 13, 0.6,
                 "이미 개발된 Fryndo MVP 보유  |  빠른 실행력  |  린 스타트업 방식",
                 font_size=18, color=DARK, align=PP_ALIGN.CENTER)


def slide_17_roadmap(prs):
    """17. 로드맵"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, 1.2)
    add_text_box(slide, 0.8, 0.35, 14, 0.7, "1년 로드맵", font_size=40, bold=True, color=WHITE)

    quarters = [
        ("Q1", "1~3월", "기반 구축", ["AR 모듈 개발", "API 연동", "기존 앱 통합"], PRIMARY),
        ("Q2", "4~6월", "MVP 개발", ["핵심 3기능 구현", "내부 테스트", "버그 수정"], ACCENT2),
        ("Q3", "7~9월", "베타 테스트", ["테스터 100명", "실제 여행 테스트", "피드백 반영"], ACCENT),
        ("Q4", "10~12월", "정식 출시", ["앱스토어 출시", "1,000명 확보", "BM 검증"], PRIMARY_DARK),
    ]

    for i, (q, period, title, items, color) in enumerate(quarters):
        x = 0.8 + i * 3.8
        header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.8), Inches(3.5), Inches(0.8))
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.fill.background()
        add_text_box(slide, x, 1.9, 3.5, 0.6, f"{q} ({period})", font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, x, 2.8, 3.5, 0.5, title, font_size=20, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            add_text_box(slide, x + 0.2, 3.4 + j * 0.5, 3.2, 0.4, f"• {item}", font_size=14, color=GRAY)

    # 마일스톤
    add_text_box(slide, 0.8, 5.5, 14, 0.5, "🎯 주요 마일스톤", font_size=22, bold=True, color=DARK)
    milestones = [
        ("3월", "AR 프로토타입"),
        ("6월", "MVP 테스트 완료"),
        ("9월", "베타 100명"),
        ("12월", "정식 출시"),
    ]
    for i, (month, desc) in enumerate(milestones):
        x = 1 + i * 3.7
        add_text_box(slide, x, 6.1, 3.5, 0.4, month, font_size=18, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
        add_text_box(slide, x, 6.6, 3.5, 0.5, desc, font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

    add_text_box(slide, 0.8, 7.5, 14, 0.5,
                 "2년차: AR 안경 연동 (Meta Ray-Ban)  |  B2B 확장  |  글로벌 진출",
                 font_size=16, color=GRAY, align=PP_ALIGN.CENTER)


def slide_18_funding(prs):
    """18. 자금 계획"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, 1.2)
    add_text_box(slide, 0.8, 0.35, 14, 0.7, "자금 계획 (예비창업패키지 1억원)", font_size=40, bold=True, color=WHITE)

    items = [
        ("인건비", "40%", "4,000만원", "대표 인건비, 아르바이트", PRIMARY),
        ("외주 개발비", "25%", "2,500만원", "AR 모듈, UI/UX, QA", ACCENT2),
        ("API 비용", "10%", "1,000만원", "번역/음성인식 API", ACCENT),
        ("마케팅비", "15%", "1,500만원", "SNS 광고, 인플루언서", RGBColor(156, 39, 176)),
        ("운영비", "10%", "1,000만원", "서버, 기타 경비", GRAY),
    ]

    for i, (name, percent, amount, desc, color) in enumerate(items):
        y = 1.8 + i * 1.05
        add_text_box(slide, 1, y, 2.5, 0.5, name, font_size=18, bold=True, color=DARK)
        bar_width = float(percent.replace('%', '')) / 100 * 8
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(y + 0.1), Inches(bar_width), Inches(0.4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_text_box(slide, 12, y, 1.5, 0.5, percent, font_size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text_box(slide, 13.5, y, 2, 0.5, amount, font_size=14, color=GRAY)
        add_text_box(slide, 3.5, y + 0.45, 8, 0.4, desc, font_size=11, color=GRAY)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(7.2), Inches(14), Inches(1))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.fill.background()
    add_text_box(slide, 1.5, 7.4, 13, 0.6,
                 "💰 자금 조달: 예비창업패키지 1억  →  시드 투자 3억 (MVP 검증 후)  →  시리즈 A (PMF 달성 후)",
                 font_size=16, color=DARK, align=PP_ALIGN.CENTER)


def slide_19_risk(prs):
    """19. 리스크 & 대응 방안 (NEW)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide, 1.2)
    add_text_box(slide, 0.8, 0.35, 14, 0.7, "리스크 & 대응 방안", font_size=40, bold=True, color=WHITE)

    risks = [
        ("🔧 기술 리스크", "AR 기능 개발 지연", "외주 개발자 사전 확보\n스마트폰 AR 우선 (안경은 향후)", RGBColor(244, 67, 54)),
        ("📈 시장 리스크", "사용자 확보 어려움", "기존 여행 커뮤니티 제휴\n인플루언서 마케팅 집중", RGBColor(255, 152, 0)),
        ("🏢 경쟁 리스크", "대기업 유사 서비스 출시", "그룹 여행 니치 시장 집중\n빠른 실행으로 선점", RGBColor(33, 150, 243)),
        ("💰 자금 리스크", "예창패 미선정 시", "자체 자금으로 최소 MVP 완성\n다른 정부지원 재도전", RGBColor(156, 39, 176)),
    ]

    for i, (icon_title, risk, response, color) in enumerate(risks):
        y = 1.8 + i * 1.4

        # 리스크 박스
        risk_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(4.5), Inches(1.2))
        risk_box.fill.solid()
        risk_box.fill.fore_color.rgb = RGBColor(255, 235, 238)
        risk_box.line.color.rgb = color

        add_text_box(slide, 1, y + 0.1, 4.2, 0.4, icon_title, font_size=16, bold=True, color=color)
        add_text_box(slide, 1, y + 0.55, 4.2, 0.6, risk, font_size=14, color=DARK)

        # 화살표
        add_text_box(slide, 5.5, y + 0.35, 0.8, 0.5, "→", font_size=24, color=GRAY, align=PP_ALIGN.CENTER)

        # 대응 박스
        resp_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(y), Inches(8.5), Inches(1.2))
        resp_box.fill.solid()
        resp_box.fill.fore_color.rgb = RGBColor(232, 245, 233)
        resp_box.line.color.rgb = ACCENT2

        add_text_box(slide, 6.7, y + 0.1, 8, 0.4, "✅ 대응 방안", font_size=14, bold=True, color=ACCENT2)
        add_text_box(slide, 6.7, y + 0.5, 8, 0.7, response, font_size=13, color=DARK)

    # 하단 메시지
    add_text_box(slide, 0.8, 7.5, 14, 0.5,
                 "💡 리스크를 인지하고 선제적으로 대응 방안을 준비하고 있습니다",
                 font_size=16, color=PRIMARY, align=PP_ALIGN.CENTER)


def slide_20_closing(prs):
    """20. 마무리"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, PRIMARY)
    add_text_box(slide, 1, 2.0, 14, 0.5, "Our Vision", font_size=24, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 2.8, 14, 1.5,
                 '"언어가 달라도,\n전 세계 누구와도\n함께 여행하는 세상"',
                 font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(5.5), Inches(4), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = WHITE
    line.line.fill.background()

    add_text_box(slide, 1, 6.0, 14, 0.8, "Fryndo AR Companion",
                 font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 7.2, 14, 0.6, "감사합니다",
                 font_size=28, color=WHITE, align=PP_ALIGN.CENTER)


def slide_21_qna(prs):
    """21. Q&A"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK)
    add_text_box(slide, 1, 3, 14, 1.2, "Q & A",
                 font_size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 5, 14, 0.6, "질문 있으시면 말씀해 주세요",
                 font_size=24, color=GRAY, align=PP_ALIGN.CENTER)

    # [수정 필요] 실제 연락처로 변경
    add_text_box(slide, 1, 7, 14, 0.5, "📧 [이메일]  |  📱 [전화번호]",
                 font_size=18, color=GRAY, align=PP_ALIGN.CENTER)


# ========== 슬라이드 생성 실행 ==========

slide_01_cover(prs)
slide_02_toc(prs)
slide_03_problem1(prs)
slide_04_problem2(prs)
slide_05_solution(prs)
slide_06_feature1(prs)
slide_07_feature2(prs)
slide_08_feature3(prs)
slide_09_scenario(prs)
slide_10_demo(prs)
slide_11_mvp(prs)        # NEW: 기존 MVP 현황
slide_12_tech(prs)
slide_13_market(prs)     # 출처 추가
slide_14_competitor(prs) # 경쟁사 확대
slide_15_business(prs)
slide_16_team(prs)       # 경력 수정
slide_17_roadmap(prs)
slide_18_funding(prs)
slide_19_risk(prs)       # NEW: 리스크 대응
slide_20_closing(prs)
slide_21_qna(prs)

# 저장
output_path = r"C:\Develop\workspace\00.Down_ai\docs\Fryndo_AR_Companion_IR.pptx"
prs.save(output_path)
print(f"PPT 생성 완료: {output_path}")
print(f"총 {len(prs.slides)}개 슬라이드")
